#!/usr/bin/env python3
"""
Generate an editable PowerPoint deck for the OTP Messenger presentation.

Requires: python-pptx   (pip install python-pptx)
Run:      python3 presentation/build_deck.py
Output:   presentation/OTP_Messenger_Deck.pptx

The generated .pptx is fully editable — open it in PowerPoint / Keynote / Google Slides,
add your names on slide 1, and paste demo screenshots onto slides 2, 7 and 15.
Slide content mirrors 05_slide_deck.md; speaker notes point back to 01_presentation_script.md.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x12, 0x26, 0x3A)
ACCENT = RGBColor(0x2E, 0x9C, 0xCA)
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)
DARK   = RGBColor(0x22, 0x2A, 0x33)
GREY   = RGBColor(0x5A, 0x67, 0x72)
GREYL  = RGBColor(0xB9, 0xC4, 0xCE)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x1E, 0x1E, 0x1E)
CODEFG = RGBColor(0xE6, 0xE6, 0xE6)
FONT   = "Calibri"
MONO   = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── low-level helpers ─────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, l, t, w, h):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)).text_frame
    tf.word_wrap = True
    return tf


def _p(tf, text, size, color, bold=False, level=0, first=False,
       align=PP_ALIGN.LEFT, after=6, font=FONT):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.level = level
    p.alignment = align
    p.space_after = Pt(after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return p


def _bar(slide, l, t, w=3.2, h=0.06, color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _footer(slide, n):
    _p(_box(slide, 0.5, 7.02, 10.0, 0.4), "OTP Messenger · UFCE87-15-3", 10, GREY, first=True)
    _p(_box(slide, 12.0, 7.02, 0.9, 0.4), str(n), 10, GREY, first=True, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _base(title):
    slide = prs.slides.add_slide(BLANK)
    _bg(slide, WHITE)
    _p(_box(slide, 0.6, 0.42, 12.1, 1.0), title, 30, NAVY, bold=True, first=True)
    _bar(slide, 0.62, 1.32)
    return slide


# ── slide builders ────────────────────────────────────────────────────────────
def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    tf = _box(s, 0.9, 2.05, 11.5, 2.0)
    _p(tf, "OTP Messenger", 52, WHITE, bold=True, first=True, after=2)
    _p(tf, "Security Assessment & Hardening", 30, ACCENT, bold=True)
    _bar(s, 0.95, 4.15, w=4.5)
    tf2 = _box(s, 0.9, 4.45, 11.5, 2.2)
    _p(tf2, "Cyber Security Engineering (UFCE87-15-3) · Recorded Video (Resit)",
       18, WHITE, first=True, after=16)
    _p(tf2, "[Student 1 name — student number]     ·     [Student 2 name — student number]",
       18, WHITE, after=12)
    _p(tf2, "July 2026", 15, GREYL)
    _notes(s, "Speaker A · ~0:00. Introduce both of you + one line on what the app is. ~20s.")
    return s


def section_slide(title, n):
    s = prs.slides.add_slide(BLANK)
    _bg(s, NAVY)
    _p(_box(s, 0.9, 3.1, 11.5, 1.6), title, 44, WHITE, bold=True, first=True)
    _bar(s, 0.95, 4.55, w=4.0)
    _footer(s, n)
    return s


def bullets_slide(title, items, n, demo=None, note=""):
    s = _base(title)
    body = _box(s, 0.8, 1.7, 11.8, 4.9)
    for i, it in enumerate(items):
        text, lvl = it if isinstance(it, tuple) else (it, 0)
        glyph = "•  " if lvl == 0 else "–  "
        _p(body, glyph + text, 20 if lvl == 0 else 18,
           DARK if lvl == 0 else GREY, level=lvl, first=(i == 0), after=11)
    if demo:
        _p(_box(s, 0.8, 6.45, 11.8, 0.5), "▶  " + demo, 16, ACCENT, bold=True, first=True)
    _footer(s, n)
    if note:
        _notes(s, note)
    return s


def columns_slide(title, columns, n, footer_note=None, note=""):
    s = _base(title)
    gap, total_w, left0, top = 0.35, 12.0, 0.7, 1.8
    col_w = (total_w - gap * (len(columns) - 1)) / len(columns)
    for i, (heading, items) in enumerate(columns):
        l = left0 + i * (col_w + gap)
        _p(_box(s, l, top, col_w, 0.6), heading, 18, NAVY, bold=True, first=True)
        _bar(s, l + 0.02, top + 0.55, w=col_w - 0.25, h=0.045)
        ib = _box(s, l, top + 0.75, col_w, 4.2)
        for j, it in enumerate(items):
            _p(ib, "•  " + it, 14, DARK, first=(j == 0), after=8)
    if footer_note:
        _p(_box(s, 0.7, 6.35, 12.0, 0.6), footer_note, 14, GREY, first=True)
    _footer(s, n)
    if note:
        _notes(s, note)
    return s


def table_slide(title, headers, rows, n, note="", verdict=None, bold_last=True):
    s = _base(title)
    nrows, ncols = len(rows) + 1, len(headers)
    tbl = s.shapes.add_table(nrows, ncols, Inches(0.8), Inches(1.85),
                             Inches(11.7), Inches(0.45 * nrows)).table
    for c, h in enumerate(headers):
        _cell(tbl.cell(0, c), h, WHITE, NAVY, bold=True, size=16)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            bold = bold_last and c == ncols - 1
            _cell(tbl.cell(r, c), val, DARK, LIGHT if r % 2 else WHITE, bold=bold, size=14)
    if verdict:
        _p(_box(s, 0.8, 1.9 + 0.45 * nrows + 0.15, 11.7, 0.6), verdict, 16, NAVY,
           bold=True, first=True)
    _footer(s, n)
    if note:
        _notes(s, note)
    return s


def _cell(cell, text, fg, bg, bold=False, size=14, align=PP_ALIGN.LEFT):
    cell.fill.solid()
    cell.fill.fore_color.rgb = bg
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = cell.margin_right = Inches(0.1)
    cell.margin_top = cell.margin_bottom = Inches(0.03)
    cell.text = text
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = fg
            r.font.name = FONT


def findings_slide(n):
    s = _base("Three headline findings")
    body = _box(s, 0.8, 1.65, 11.8, 2.1)
    items = [
        "Debug mode ON → interactive debugger = remote code execution  (CAT I)",
        "HTTP only → OTP key in cleartext + plaintext written to logs  (CAT I / CAT II)",
        "Truncating XOR: zip() stops at the shorter input → bytes left unencrypted",
    ]
    for i, t in enumerate(items):
        _p(body, "•  " + t, 18, DARK, first=(i == 0), after=11)
    code = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(4.05),
                              Inches(11.8), Inches(1.55))
    code.fill.solid()
    code.fill.fore_color.rgb = CODEBG
    code.line.fill.background()
    code.shadow.inherit = False
    ctf = code.text_frame
    ctf.word_wrap = True
    ctf.margin_left = ctf.margin_right = Inches(0.2)
    ctf.margin_top = ctf.margin_bottom = Inches(0.12)
    lines = [
        "data = plaintext.encode('utf-8')",
        "key  = os.urandom(len(data))",
        "ciphertext = bytes(p ^ k for p, k in zip(data, key))   # zip stops at shorter → silent truncation",
    ]
    for i, ln in enumerate(lines):
        _p(ctf, ln, 13, CODEFG, first=(i == 0), after=2, font=MONO)
    _p(_box(s, 0.8, 5.75, 11.8, 0.5),
       "▶  Live demo: debugger banner · cleartext key over HTTP · plaintext in the log",
       15, ACCENT, bold=True, first=True)
    _footer(s, n)
    _notes(s, "Speaker A · ~3:00. Runbook O1/O2/O3. Stress the XOR bug was found by reading code.")
    return s


# ── build ─────────────────────────────────────────────────────────────────────
def build():
    title_slide()                                                                    # 1
    bullets_slide("What is the OTP Messenger?", [
        "Browser-based one-time pad encryption tool",
        "Type a message → get ciphertext + a single-use key",
        "Key shown once, never stored; sent out-of-band",
        "Browser → HTTPS → Nginx → Flask/Gunicorn",
    ], 2, demo="Live demo: encrypt “MEETING AT 0900” → decrypt it back",
        note="Speaker A · ~0:20. Plain-English framing, cut to browser demo (runbook F1).")   # 2
    bullets_slide("Our approach", [
        "We assess — we didn't build it",
        "Four areas: Vulnerabilities · Compliance · Changes · Testing",
        "Baseline vs Hardened — compared throughout",
        "Split: A = app & fixes · B = standards & assurance",
    ], 3, note="Speaker A · ~1:20. Set expectations + hand-off plan.")               # 3
    section_slide("Identified Vulnerabilities", 4)                                    # 4
    bullets_slide("Methodology & tools", [
        "Static code review — Python, Dockerfiles, Nginx, Compose",
        "Tooling — Trivy (image CVEs) · Nmap (services) · fuzzing",
        "Dynamic testing on the running app",
        "Malicious-code review → none found (issues are vulns / misconfig)",
    ], 5, note="Speaker A · ~2:00. Say the malicious-code check out loud — the brief asks for it.")  # 5
    bullets_slide("Findings at a glance", [
        "5 × CAT I  (Critical)",
        "9 × CAT II  (High)",
        "3 × CAT III  (Low)",
        "= 17 findings — baseline NOT authorised to deploy",
    ], 6, note="Speaker A · ~2:30. Scale of the problem.")                            # 6
    findings_slide(7)                                                                 # 7
    section_slide("Required Compliance Standards", 8)                                 # 8
    bullets_slide("Standards selected — and why", [
        "ASD STIG V5R3 — it's a web application",
        "NIST SP 800-52 Rev 2 — it terminates TLS",
        "Container Platform SRG — it runs in containers",
        "Web Server SRG — Nginx ingress",
        "Found via the DoD STIG library, searched on functionality",
    ], 9, note="Speaker B · ~5:00. Explain the selection reasoning, not just the list.")  # 9
    table_slide("Baseline shortcomings",
                ["Finding", "Control", "Standard"],
                [["Debug mode enabled", "APSC-DV-002530", "ASD STIG"],
                 ["Hardcoded secret key", "APSC-DV-003280", "ASD STIG"],
                 ["No HTTPS", "APSC-DV-002000", "ASD STIG / NIST 800-52"],
                 ["Plaintext in logs", "APSC-DV-003200", "ASD STIG"],
                 ["Root container", "V-205072", "Container SRG"],
                 ["No dropped capabilities", "V-205070", "Container SRG"],
                 ["End-of-life images", "V-205076", "Container SRG"]],
                10, verdict="Verdict: not authorised against any standard.",
                bold_last=False,
                note="Speaker B · ~6:00. Map concrete failures to real control IDs.")  # 10
    table_slide("Accepted risks — all CAT III",
                ["Risk", "Justification", "Fix before production"],
                [["Self-signed certificate", "No DoD PKI in dev; RSA-4096 / SHA-256 meet NIST", "Swap in CA-issued cert"],
                 ["FIPS 140-2 not validated", "Host config, not code; OS CSPRNG is FIPS-equivalent", "Deploy on FIPS-mode host"],
                 ["OTP key distribution", "Out-of-band procedural control", "Document in key-mgmt policy"]],
                11, bold_last=False,
                note="Speaker B · ~7:15. Each accepted risk is environmental, with a timeline.")  # 11
    section_slide("Changes Made", 12)                                                # 12
    columns_slide("Remediation by layer", [
        ("Application", ["Gunicorn (debug off)", "CSPRNG secret key", "Flask-WTF CSRF",
                         "Validation + 32 KB cap", "Metadata-only logging", "XOR length check"]),
        ("Transport", ["TLS 1.2 / 1.3", "NIST cipher suites", "HTTP→HTTPS 301",
                       "HSTS + CSP + headers"]),
        ("Container", ["Non-root user", "cap_drop: ALL", "Read-only filesystem",
                       "Resource limits", "Healthchecks", "No host port"]),
    ], 13, footer_note="Enhancements: character counter · copy-to-clipboard · /health endpoint",
        note="Speaker A · ~8:15. Fix without breaking function; walk the three layers.")  # 13
    section_slide("Testing & Verification", 14)                                      # 14
    bullets_slide("Functionality + vulnerability re-tests", [
        "Encrypt → decrypt round-trip  ✓",
        "Debug gone (Gunicorn)  ✓",
        "TLS enforced; TLS 1.1 rejected; HTTP→301  ✓",
        "Port 5000 refused  ✓",
        "CSRF → 400  ✓",
        "Oversized input → 413  ✓",
    ], 15, demo="Live demo: tick each one live (runbook F1–F6)",
        note="Speaker B · ~10:45. Each closed finding, verified on camera.")          # 15
    bullets_slide("The iterative loop", [
        "Change → Re-test → Re-scan → Re-verify",
        "First hardened scan was NOT clean:",
        ("Gunicorn 21.2.0 → CVE-2024-1135 (request smuggling)", 1),
        ("Nginx image → ~20 CVEs", 1),
        "Fixed: Gunicorn → 22.0.0 · Nginx apk upgrade",
        "Re-scanned: Gunicorn CVE gone · Nginx image → 0",
    ], 16, note="Speaker B · ~12:00. This is the iterative-process evidence the brief wants.")  # 16
    table_slide("Before vs After",
                ["Metric", "Original", "Final"],
                [["CAT I", "5", "0"],
                 ["CAT II", "9", "0"],
                 ["Image CVEs (HIGH/CRIT)", "1,403", "14"],
                 ["STIG status", "Not authorised", "Conditional"],
                 ["Core function", "Works", "Works"]],
                17, note="Speaker B · ~13:15. Headline deltas + functionality preserved.")  # 17
    bullets_slide("Conclusion", [
        "Meets baseline: 0 CAT I / 0 CAT II",
        "3 accepted risks — environmental, justified, with a fix path",
        "Functionality intact",
        "Not “perfect” — but understood & justified",
        "Thank you — questions?",
    ], 18, note="Speaker B · ~14:00. Honest close + thanks. Finish by 14:30.")         # 18
    bullets_slide("References (UWE-Harvard)", [
        "DISA (2024) Application Security and Development STIG, V5R3",
        "NIST (2019) SP 800-52 Rev. 2 — Guidelines for TLS",
        "DISA Container Platform SRG; DISA Web Server SRG",
        "Tools: Flask, Flask-WTF, Gunicorn, Nginx, Trivy (Aqua Security)",
        "Generative AI (Claude) used to structure the deck & explain STIG wording; "
        "findings and testing are our own (adapt per UWE GenAI policy)",
    ], 19, note="Format each reference fully in UWE-Harvard. Keep on screen a few seconds.")  # 19

    out = Path(__file__).resolve().parent / "OTP_Messenger_Deck.pptx"
    prs.save(str(out))
    print(f"[+] Wrote {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
